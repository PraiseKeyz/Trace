from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand colours ──────────────────────────────────────────────────────────────
DARK     = RGBColor(0x02, 0x06, 0x17)
ORANGE   = RGBColor(0xF9, 0x73, 0x16)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xF9, 0xF6, 0xF0)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
CARD_BG  = RGBColor(0x0c, 0x15, 0x2a)
DIM      = RGBColor(0x1e, 0x2d, 0x4a)
FOOTNOTE = RGBColor(0x44, 0x55, 0x70)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def tb(slide, text, x, y, w, h,
       size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
       italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return txb

def tb_lines(slide, lines, x, y, w, h, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = 'Calibri'
    return txb

def base(s, label):
    rect(s, 0, 0, W, H, DARK)
    tb(s, label, Inches(0.6), Inches(0.45), Inches(8), Inches(0.4),
       size=11, bold=True, color=ORANGE)
    rect(s, Inches(0.6), Inches(7.1), Inches(12.1), Inches(0.03), ORANGE)
    ring = s.shapes.add_shape(9, Inches(11.5), Inches(0.2), Inches(1.7), Inches(1.7))
    ring.fill.background()
    ring.line.color.rgb = DIM
    ring.line.width = Pt(16)

def card_base(s, x, y, w, h):
    rect(s, x, y, w, h, CARD_BG)
    rect(s, x, y, w, Inches(0.05), ORANGE)

def footnote(s, text):
    tb(s, text, Inches(0.6), Inches(7.2), Inches(12), Inches(0.3),
       size=8, color=FOOTNOTE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
rect(s1, 0, 0, W, H, DARK)
rect(s1, Inches(0.6), Inches(1.6), Inches(0.06), Inches(4.3), ORANGE)
tb(s1, 'TRACE', Inches(0.85), Inches(1.5), Inches(8), Inches(2.2),
   size=96, bold=True, color=WHITE)
tb(s1, 'Economic Identity for the Invisible',
   Inches(0.88), Inches(3.55), Inches(9), Inches(0.8),
   size=26, bold=False, color=ORANGE)
tb(s1, 'Connecting informal traders & gig workers to the economy they power.',
   Inches(0.88), Inches(4.25), Inches(9), Inches(0.7),
   size=16, bold=False, color=MUTED)
rect(s1, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.03), ORANGE)
tb(s1, 'Squad Hackathon 3.0',
   Inches(0.6), Inches(6.55), Inches(5), Inches(0.5),
   size=13, bold=False, color=MUTED)
tb(s1, 'Challenge 02  ·  Smart Systems: The Intelligent Economy',
   Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.5),
   size=13, bold=False, color=MUTED, align=PP_ALIGN.RIGHT)
ring_outer = s1.shapes.add_shape(9, Inches(9.8), Inches(0.6), Inches(3.5), Inches(3.5))
ring_outer.fill.background()
ring_outer.line.color.rgb = ORANGE
ring_outer.line.width = Pt(1.5)
ring_inner = s1.shapes.add_shape(9, Inches(10.3), Inches(1.1), Inches(2.5), Inches(2.5))
ring_inner.fill.background()
ring_inner.line.color.rgb = RGBColor(0x10, 0x18, 0x30)
ring_inner.line.width = Pt(18)
dot = s1.shapes.add_shape(9, Inches(11.1), Inches(4.5), Inches(0.25), Inches(0.25))
dot.fill.solid()
dot.fill.fore_color.rgb = ORANGE
dot.line.fill.background()


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
rect(s2, 0, 0, W, H, DARK)
tb(s2, '01  —  PROBLEM', Inches(0.6), Inches(0.45), Inches(6), Inches(0.4),
   size=11, bold=True, color=ORANGE)
tb_lines(s2, [
    ('70% Work.', 54, True, WHITE),
    ('None of it Counts.', 54, True, ORANGE),
], Inches(0.6), Inches(0.9), Inches(9), Inches(2.1))
tb(s2,
   'Across sub-Saharan Africa, millions sustain families through informal trade — '
   'yet they are invisible to banks, lenders, and the systems that could empower them.',
   Inches(0.6), Inches(2.85), Inches(8.2), Inches(1.0),
   size=15, bold=False, color=MUTED)
CARD_Y  = Inches(4.05)
CARD_H  = Inches(2.7)
CARD_W  = Inches(3.7)
GAP     = Inches(0.28)
CARD_X1 = Inches(0.6)
CARD_X2 = CARD_X1 + CARD_W + GAP
CARD_X3 = CARD_X2 + CARD_W + GAP
for cx in [CARD_X1, CARD_X2, CARD_X3]:
    rect(s2, cx, CARD_Y, CARD_W, CARD_H, CARD_BG)
    rect(s2, cx, CARD_Y, CARD_W, Inches(0.04), ORANGE)
tb(s2, '70%',  CARD_X1 + Inches(0.25), CARD_Y + Inches(0.35), Inches(3.2), Inches(0.9),
   size=52, bold=True, color=ORANGE)
tb(s2, 'of employment in sub-Saharan Africa\ncomes from the informal economy',
   CARD_X1 + Inches(0.25), CARD_Y + Inches(1.2), Inches(3.2), Inches(1.2),
   size=13, bold=False, color=MUTED)
tb(s2, '0',    CARD_X2 + Inches(0.25), CARD_Y + Inches(0.35), Inches(3.2), Inches(0.9),
   size=52, bold=True, color=ORANGE)
tb(s2, 'credit score  ·  0 loan history\n0 financial identity to show banks',
   CARD_X2 + Inches(0.25), CARD_Y + Inches(1.2), Inches(3.2), Inches(1.2),
   size=13, bold=False, color=MUTED)
tb(s2, '$17.6B', CARD_X3 + Inches(0.25), CARD_Y + Inches(0.35), Inches(3.2), Inches(0.9),
   size=40, bold=True, color=ORANGE)
tb(s2, 'in informal cross-border trade in\nSADC alone — uncounted, unserved',
   CARD_X3 + Inches(0.25), CARD_Y + Inches(1.2), Inches(3.2), Inches(1.2),
   size=13, bold=False, color=MUTED)
ring2 = s2.shapes.add_shape(9, Inches(10.5), Inches(0.3), Inches(2.6), Inches(2.6))
ring2.fill.background()
ring2.line.color.rgb = RGBColor(0x1e, 0x2d, 0x4a)
ring2.line.width = Pt(22)
tb(s2, 'Source: Global Initiative · SADC Informal Trade Report',
   Inches(0.6), Inches(7.1), Inches(8), Inches(0.35),
   size=9, bold=False, color=FOOTNOTE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — TARGET USER
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
base(s3, '02  —  TARGET USER')
tb(s3, 'Meet the People We Serve',
   Inches(0.6), Inches(0.85), Inches(10), Inches(0.75),
   size=38, bold=True)
tb(s3,
   'We spoke to informal traders and gig workers across Lagos and Abuja. Two of millions '
   'who exist outside the financial system — not by choice, but by design.',
   Inches(0.6), Inches(1.62), Inches(11.8), Inches(0.5),
   size=13, color=MUTED)

CW = Inches(5.85)
CY = Inches(2.25)
CH = Inches(3.85)
CX1 = Inches(0.6)
CX2 = Inches(7.0)
P   = Inches(0.28)

for cx in [CX1, CX2]:
    card_base(s3, cx, CY, CW, CH)

tb(s3, 'AMAKA', CX1+P, CY+Inches(0.18), CW-P*2, Inches(0.55), size=28, bold=True)
tb(s3, '34  ·  Fabric Trader  ·  Lagos Island Market',
   CX1+P, CY+Inches(0.72), CW-P*2, Inches(0.35), size=12, color=ORANGE)
rect(s3, CX1+P, CY+Inches(1.12), CW-P*2, Inches(0.015), DIM)
tb(s3, '8 yrs in business   ·   ₦45K/week revenue   ·   0 credit score',
   CX1+P, CY+Inches(1.18), CW-P*2, Inches(0.35), size=10, color=MUTED)
tb(s3,
   'Runs a fabric stall in Lagos Island Market. Eight years of trading, consistent weekly '
   'revenue, and loyal customers — but no bank has ever given her a loan. Without a payslip '
   'or credit history, she does not exist to the formal financial system.',
   CX1+P, CY+Inches(1.63), CW-P*2, Inches(1.35), size=12, color=MUTED)
tb(s3, '"I work every single day. Why can\'t the bank see that?"',
   CX1+P, CY+Inches(3.08), CW-P*2, Inches(0.68),
   size=13, bold=True, italic=True, color=ORANGE)

tb(s3, 'EMEKA', CX2+P, CY+Inches(0.18), CW-P*2, Inches(0.55), size=28, bold=True)
tb(s3, '26  ·  Gig Delivery Rider  ·  Abuja / FCT',
   CX2+P, CY+Inches(0.72), CW-P*2, Inches(0.35), size=12, color=ORANGE)
rect(s3, CX2+P, CY+Inches(1.12), CW-P*2, Inches(0.015), DIM)
tb(s3, '3 gig platforms   ·   ₦85K/month   ·   No verifiable record',
   CX2+P, CY+Inches(1.18), CW-P*2, Inches(0.35), size=10, color=MUTED)
tb(s3,
   'Works across three delivery apps simultaneously. His income is real and consistent '
   'but fragmented across platforms — no single institution can see the full picture. '
   'Applied for a business loan to buy his own motorbike. Rejected. No credit history.',
   CX2+P, CY+Inches(1.63), CW-P*2, Inches(1.35), size=12, color=MUTED)
tb(s3, '"My skills are real. I just need a way to prove them."',
   CX2+P, CY+Inches(3.08), CW-P*2, Inches(0.68),
   size=13, bold=True, italic=True, color=ORANGE)

for i, (num, label) in enumerate([
    ('39M',   'informal workers in Nigeria'),
    ('<5%',   'have access to formal credit'),
    ('₦120K', 'average monthly informal income'),
]):
    sx = Inches(0.6 + i * 4.24)
    tb(s3, num,   sx, Inches(6.22), Inches(4.0), Inches(0.52), size=26, bold=True)
    tb(s3, label, sx, Inches(6.72), Inches(4.0), Inches(0.35), size=11, color=MUTED)

footnote(s3, 'Source: ILO Informal Economy Report 2023  ·  EFInA Access to Financial Services Nigeria 2023')


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SOLUTION OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
base(s4, '03  —  SOLUTION OVERVIEW')
tb(s4, 'One Platform. Four Engines.',
   Inches(0.6), Inches(0.85), Inches(10), Inches(0.75), size=38, bold=True)
tb(s4,
   'Trace turns every transaction, gig, and peer interaction into a verified economic signal — '
   'building the digital identity that unlocks financial services for the people who have always been excluded.',
   Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5), size=13, color=MUTED)

FW = Inches(5.85)
FH = Inches(2.3)

for (fx, fy, title, body) in [
    (Inches(0.6),  Inches(2.2),
     '01  /  Digital Onboarding',
     'Phone-first, 3-minute signup. Name, skills, trade category, location, and preferred language. '
     'Data-sharing consent captured upfront. A Squad virtual account is auto-generated on completion — '
     'every user has a verified payment identity from minute one.'),
    (Inches(6.9),  Inches(2.2),
     '02  /  Economic Identity Score',
     'AI-powered 0–100 score built from four signals: transaction history (40%), peer vouches (25%), '
     'platform activity (20%), and profile completeness (15%). Calculated by a Python AI service '
     'via gRPC and updated automatically after every Squad transaction and gig completion.'),
    (Inches(0.6),  Inches(4.7),
     '03  /  Opportunity Marketplace',
     'AI matches gig workers and job seekers to opportunities using skill overlap (40%), proximity (25%), '
     'historical success rate (20%), and language match (15%). Employer payments are locked in Squad '
     'escrow before the gig starts — trust built into every transaction by design.'),
    (Inches(6.9),  Inches(4.7),
     '04  /  Finance Gateway',
     'Score unlocks partner financial products: micro-loans (₦50K–₦500K), savings plans, and '
     'business insurance — offered by licensed microfinance institutions who query the Trace Partner API. '
     'Approved funds are disbursed directly to the user\'s Squad virtual account.'),
]:
    card_base(s4, fx, fy, FW, FH)
    tb(s4, title, fx+Inches(0.25), fy+Inches(0.15), FW-Inches(0.4), Inches(0.4),
       size=12, bold=True, color=ORANGE)
    tb(s4, body,  fx+Inches(0.25), fy+Inches(0.55), FW-Inches(0.4), FH-Inches(0.65),
       size=12, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SQUAD API INTEGRATION
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
base(s5, '04  —  SQUAD API INTEGRATION')
tb(s5, "Squad Isn't a Feature. It's the Foundation.",
   Inches(0.6), Inches(0.85), Inches(11), Inches(0.75), size=36, bold=True)
tb(s5,
   'Every critical flow on Trace runs through Squad. Remove Squad, and the platform loses '
   'its ability to verify identity, capture economic signals, and disburse financial products.',
   Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5), size=13, color=MUTED)

AW = Inches(3.75)
AH = Inches(3.8)
AY = Inches(2.2)

for (ax, num, title, body, endpoint) in [
    (Inches(0.6), '01', 'Virtual Accounts',
     'Every user receives a Squad-powered Wema Bank virtual account the moment they complete '
     'onboarding. This is their economic home — where payments arrive, escrow funds are released, '
     'and loan disbursements land. All transaction history is captured here.',
     'POST /squad/my-virtual-account'),
    (Inches(4.75), '02', 'Escrow & Transfers',
     'When a gig is accepted, the employer\'s payment is locked in Squad escrow before the '
     'worker starts. Trace calls Squad to release only on confirmed completion — zero trust '
     'required between two strangers meeting for the first time.',
     'POST /squad/escrow/:id/lock\nPOST /squad/escrow/:id/release'),
    (Inches(8.9), '03', 'Webhooks & Verification',
     'Squad fires a signed webhook on every completed payment. Trace validates the signature, '
     'records the transaction to the database, and queues a score recalculation via BullMQ '
     'into the Python gRPC service. Every payment makes the user more creditworthy.',
     'POST /webhooks/squad\nGET  /squad/transactions/verify/:ref'),
]:
    card_base(s5, ax, AY, AW, AH)
    tb(s5, num,   ax+Inches(0.25), AY+Inches(0.18), AW-Inches(0.4), Inches(0.5),
       size=28, bold=True, color=ORANGE)
    tb(s5, title, ax+Inches(0.25), AY+Inches(0.68), AW-Inches(0.4), Inches(0.4),
       size=16, bold=True)
    tb(s5, body,  ax+Inches(0.25), AY+Inches(1.12), AW-Inches(0.4), Inches(1.82),
       size=12, color=MUTED)
    rect(s5, ax+Inches(0.25), AY+AH-Inches(0.68), AW-Inches(0.4), Inches(0.58),
         RGBColor(0x06, 0x0e, 0x1f))
    tb(s5, endpoint, ax+Inches(0.35), AY+AH-Inches(0.63), AW-Inches(0.55), Inches(0.53),
       size=9, bold=True, color=ORANGE)

rect(s5, Inches(0.6),  Inches(6.2), Inches(12.1), Inches(0.72), RGBColor(0x0a, 0x18, 0x38))
rect(s5, Inches(0.6),  Inches(6.2), Inches(0.05),  Inches(0.72), ORANGE)
tb(s5,
   'Judges\' note: Squad powers virtual account creation, escrow lock/release, webhook-driven '
   'score recalculation, and financial product disbursement. This is not a token integration — '
   'Squad is the verified data spine of Trace.',
   Inches(0.85), Inches(6.27), Inches(11.6), Inches(0.6), size=12)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — AI & DATA INTELLIGENCE
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
base(s6, '05  —  AI & DATA INTELLIGENCE')
tb(s6, 'A Score That Grows With You',
   Inches(0.6), Inches(0.85), Inches(9), Inches(0.75), size=38, bold=True)
tb(s6,
   'The Economic Identity Score is calculated by a Python AI microservice via gRPC — never '
   'manually set. It updates automatically after every Squad transaction, gig, and peer vouch.',
   Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5), size=13, color=MUTED)

# Left: Score display card
card_base(s6, Inches(0.6), Inches(2.15), Inches(5.8), Inches(4.75))
tb(s6, 'IDENTITY SCORE',
   Inches(0.9), Inches(2.35), Inches(5.2), Inches(0.38), size=11, bold=True, color=ORANGE)
tb(s6, '73', Inches(0.9), Inches(2.68), Inches(2.8), Inches(1.35), size=88, bold=True)
tb(s6, '/ 100', Inches(2.95), Inches(3.3), Inches(1.5), Inches(0.55), size=18, color=MUTED)

rect(s6, Inches(0.9), Inches(4.15), Inches(1.6), Inches(0.38), RGBColor(0x05, 0x3a, 0x21))
tb(s6, 'LOW RISK', Inches(0.9), Inches(4.15), Inches(1.6), Inches(0.38),
   size=11, bold=True, color=RGBColor(0x34, 0xd3, 0x99), align=PP_ALIGN.CENTER)
tb(s6, '✓  Finance eligible up to ₦200,000',
   Inches(0.9), Inches(4.65), Inches(5.0), Inches(0.38),
   size=12, color=RGBColor(0x34, 0xd3, 0x99))

rect(s6, Inches(0.9), Inches(5.15), Inches(5.2), Inches(0.015), DIM)
tb_lines(s6, [
    ('Score formula:', 10, True, MUTED),
    ('score = tx_history(40%) + vouches(25%) + activity(20%) + profile(15%)', 10, False, MUTED),
], Inches(0.9), Inches(5.22), Inches(5.2), Inches(0.55))
tb(s6, 'Risk tiers:  High (0–29)  ·  Medium (30–54)  ·  Low (55–74)  ·  Very Low (75–100)',
   Inches(0.9), Inches(5.85), Inches(5.2), Inches(0.38), size=10, color=MUTED)
tb(s6, 'Python FastAPI  ·  gRPC  ·  scikit-learn  ·  BullMQ async queue',
   Inches(0.9), Inches(6.38), Inches(5.2), Inches(0.38), size=9, color=FOOTNOTE)

# Right: Signal bars
RX = Inches(7.0)
BW = Inches(5.9)
tb(s6, 'HOW THE SCORE IS BUILT',
   RX, Inches(2.2), BW, Inches(0.38), size=11, bold=True, color=ORANGE)

for i, (label, pct, weight) in enumerate([
    ('Transaction History',   '40%', 0.40),
    ('Community Vouching',    '25%', 0.25),
    ('Platform Activity',     '20%', 0.20),
    ('Profile Completeness',  '15%', 0.15),
]):
    by = Inches(2.72 + i * 0.9)
    tb(s6, label, RX, by, Inches(4.2), Inches(0.38), size=13)
    tb(s6, pct,   RX, by, BW,          Inches(0.38), size=13, bold=True, color=ORANGE,
       align=PP_ALIGN.RIGHT)
    rect(s6, RX, by+Inches(0.42), BW, Inches(0.2), DIM)
    rect(s6, RX, by+Inches(0.42), BW * weight, Inches(0.2), ORANGE)

rect(s6, RX, Inches(6.38), BW, Inches(0.015), DIM)
tb(s6, 'JOB MATCHING WEIGHTS',
   RX, Inches(6.42), BW, Inches(0.32), size=10, bold=True, color=MUTED)
tb(s6,
   'Skill overlap (40%)  ·  Proximity (25%)  ·  Success history (20%)  ·  Language match (15%)',
   RX, Inches(6.72), BW, Inches(0.35), size=10, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — USER FLOW
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank)
base(s7, '06  —  USER FLOW')
tb(s7, 'Invisible to Empowered in 4 Steps',
   Inches(0.6), Inches(0.85), Inches(11), Inches(0.75), size=38, bold=True)
tb(s7,
   'One continuous journey from first phone login to first financial product. '
   'Every step adds data. Every data point builds identity. Every milestone opens a new door.',
   Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5), size=13, color=MUTED)

SW  = Inches(2.7)
SH  = Inches(4.5)
SY  = Inches(2.15)
SX1 = Inches(0.75)
SX2 = SX1 + SW + Inches(0.35)
SX3 = SX2 + SW + Inches(0.35)
SX4 = SX3 + SW + Inches(0.35)

for (sx, num, title, sub, body) in [
    (SX1, '01', 'ONBOARD', '3 min  ·  phone-first',
     'Sign up with phone number. Set name, skills, trade category, preferred language, and location. '
     'Consent to data sharing. Squad virtual account auto-generated on completion — '
     'payment identity established from minute one.'),
    (SX2, '02', 'BUILD YOUR SCORE', 'Day 1 → Day 90',
     'Complete gigs through the Opportunity Marketplace. Record Squad transactions. '
     'Get vouched by peers who have transacted with you on Trace. Python AI service '
     'recalculates your Identity Score via gRPC after every event.'),
    (SX3, '03', 'GET MATCHED', 'AI-powered  ·  location-aware',
     'AI ranks open opportunities by skill overlap (40%), proximity (25%), historical '
     'success rate (20%), and language match (15%). Apply in one tap. Employer payment '
     'locked in Squad escrow before you start — guaranteed.'),
    (SX4, '04', 'ACCESS FINANCE', 'Score ≥ 55 unlocks products',
     'Your verified credit profile is exposed to partner financial institutions via the '
     'Trace Partner API. Micro-loans (₦50K–₦500K), savings plans, insurance. '
     'Funds disbursed directly to your Squad virtual account.'),
]:
    card_base(s7, sx, SY, SW, SH)
    tb(s7, num,   sx+Inches(0.2), SY+Inches(0.18), SW-Inches(0.3), Inches(0.52),
       size=30, bold=True, color=ORANGE)
    tb(s7, title, sx+Inches(0.2), SY+Inches(0.72), SW-Inches(0.3), Inches(0.42),
       size=14, bold=True)
    tb(s7, sub,   sx+Inches(0.2), SY+Inches(1.14), SW-Inches(0.3), Inches(0.32),
       size=10, color=ORANGE)
    rect(s7, sx+Inches(0.2), SY+Inches(1.5), SW-Inches(0.3), Inches(0.015), DIM)
    tb(s7, body,  sx+Inches(0.2), SY+Inches(1.58), SW-Inches(0.3), SH-Inches(1.72),
       size=11, color=MUTED)

for ax in [SX1+SW, SX2+SW, SX3+SW]:
    tb(s7, '→', ax+Inches(0.06), SY+Inches(1.9), Inches(0.28), Inches(0.45),
       size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — IMPACT POTENTIAL
# ══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank)
base(s8, '07  —  IMPACT POTENTIAL')
tb(s8, 'The Market Is Ready. The Need Is Now.',
   Inches(0.6), Inches(0.85), Inches(11), Inches(0.75), size=38, bold=True)
tb(s8,
   'Trace targets the largest underserved population in the Nigerian economy. '
   'The technology exists. The data layer is missing. We are building it.',
   Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5), size=13, color=MUTED)

IW = Inches(3.75)
IH = Inches(2.3)

for (ix, iy, num, label) in [
    (Inches(0.6),  Inches(2.15), '39M',    'informal workers in Nigeria — the primary Trace audience'),
    (Inches(4.75), Inches(2.15), '70%',    'of sub-Saharan African employment comes from the informal economy'),
    (Inches(8.9),  Inches(2.15), '$17.6B', 'in informal cross-border trade in SADC — all outside formal systems'),
    (Inches(0.6),  Inches(4.65), '10,000', 'users — Year 1 target across Lagos, Abuja & Port Harcourt'),
    (Inches(4.75), Inches(4.65), '500,000','users — Year 3 target with national rollout to all 36 states'),
    (Inches(8.9),  Inches(4.65), '₦200K',  'average first loan unlocked at Low Risk tier (score ≥ 55)'),
]:
    card_base(s8, ix, iy, IW, IH)
    nsize = 40 if len(num) <= 5 else 30
    tb(s8, num,   ix+Inches(0.25), iy+Inches(0.15), IW-Inches(0.35), Inches(0.88),
       size=nsize, bold=True, color=ORANGE)
    tb(s8, label, ix+Inches(0.25), iy+Inches(1.05), IW-Inches(0.35), IH-Inches(1.15),
       size=12, color=MUTED)

footnote(s8, 'Source: ILO 2023  ·  EFInA 2023  ·  SADC Trade Report  ·  NBS Nigeria Labour Force Survey')


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — SCALABILITY & BUSINESS MODEL
# ══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank)
base(s9, '08  —  SCALABILITY & BUSINESS MODEL')
tb(s9, 'Designed to Scale. Built to Last.',
   Inches(0.6), Inches(0.85), Inches(10), Inches(0.75), size=38, bold=True)
tb(s9,
   'Trace is not a hackathon project. It is a production-grade infrastructure play '
   'with a defensible data moat that deepens as more users join the network.',
   Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5), size=13, color=MUTED)

LW = Inches(5.7)
LX = Inches(0.6)
tb(s9, 'REVENUE STREAMS', LX, Inches(2.2), LW, Inches(0.38),
   size=11, bold=True, color=ORANGE)

for i, (rtitle, rbody) in enumerate([
    ('Transaction Fees',
     '1% on all Squad-powered transfers. 1,000 users × 5 tx/month × ₦5K avg = ₦250K/month at launch.'),
    ('Loan Origination Fee',
     '2–3% on micro-loans facilitated via partner institutions — charged to the lender, not the user.'),
    ('Premium Employer Accounts',
     '₦15,000/month for verified employers with priority listing and applicant analytics.'),
    ('Data Intelligence Reports',
     'Anonymized market demand signals sold to NGOs, government, and development finance institutions.'),
]):
    ry = Inches(2.72 + i * 0.98)
    rect(s9, LX, ry, LW, Inches(0.88), CARD_BG)
    rect(s9, LX, ry, Inches(0.04), Inches(0.88), ORANGE)
    tb(s9, rtitle, LX+Inches(0.2), ry+Inches(0.07), LW-Inches(0.28), Inches(0.35),
       size=12, bold=True)
    tb(s9, rbody,  LX+Inches(0.2), ry+Inches(0.44), LW-Inches(0.28), Inches(0.42),
       size=10, color=MUTED)

RX = Inches(6.9)
RW = Inches(5.8)
tb(s9, 'ROLLOUT ROADMAP', RX, Inches(2.2), RW, Inches(0.38),
   size=11, bold=True, color=ORANGE)

for i, (ptitle, pbody) in enumerate([
    ('Phase 1  (0–6 months)',
     'Lagos, Abuja, Port Harcourt  ·  Target: 10,000 users\nFocus: onboarding, gig matching, virtual account adoption'),
    ('Phase 2  (6–18 months)',
     'All 36 Nigerian states  ·  Target: 100,000 users\nFocus: Finance Gateway, partner institution onboarding'),
    ('Phase 3  (18–36 months)',
     'Ghana, Kenya, Rwanda  ·  Target: 500,000 users\nFocus: national credit scoring, cross-border trade intelligence'),
]):
    py = Inches(2.72 + i * 0.98)
    rect(s9, RX, py, RW, Inches(0.88), CARD_BG)
    rect(s9, RX, py, Inches(0.04), Inches(0.88), ORANGE)
    tb(s9, ptitle, RX+Inches(0.2), py+Inches(0.07), RW-Inches(0.28), Inches(0.35),
       size=12, bold=True)
    tb(s9, pbody,  RX+Inches(0.2), py+Inches(0.44), RW-Inches(0.28), Inches(0.42),
       size=10, color=MUTED)

rect(s9, RX, Inches(5.7), RW, Inches(1.1), RGBColor(0x06, 0x0e, 0x1f))
rect(s9, RX, Inches(5.7), RW, Inches(0.04), ORANGE)
tb(s9, 'TECHNOLOGY FOUNDATION', RX+Inches(0.2), Inches(5.76), RW-Inches(0.28), Inches(0.35),
   size=10, bold=True, color=ORANGE)
tb(s9,
   'NestJS  ·  Next.js (PWA)  ·  Python FastAPI  ·  gRPC  ·  PostgreSQL  ·  Redis + BullMQ  ·  Docker  ·  Squad API\n'
   'Containerized microservices — horizontal scaling from day one. No vendor lock-in.',
   RX+Inches(0.2), Inches(6.12), RW-Inches(0.28), Inches(0.65), size=10, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — RESEARCH & VALIDATION
# ══════════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(blank)
base(s10, '09  —  RESEARCH & VALIDATION')
tb(s10, 'Evidence-Backed, User-Tested.',
   Inches(0.6), Inches(0.85), Inches(10), Inches(0.75), size=38, bold=True)
tb(s10,
   'The problem is documented at scale. The market is measurable. The user demand is direct and specific.',
   Inches(0.6), Inches(1.6), Inches(10), Inches(0.45), size=13, color=MUTED)

VW = Inches(2.85)
VH = Inches(1.7)
VY = Inches(2.1)

for (vx, num, label, source) in [
    (Inches(0.6),  '70%',    'of sub-Saharan African employment is in the informal sector',          'ILO 2023'),
    (Inches(3.8),  '37%',    'of Nigerian adults have a formal bank account',                        'EFInA 2023'),
    (Inches(7.0),  '85%',    'of SME loan applications in Nigeria rejected — no credit history',     'CBN 2022'),
    (Inches(10.2), '$17.6B', 'in informal cross-border trade in SADC annually',                      'SADC Report'),
]:
    card_base(s10, vx, VY, VW, VH)
    nsize = 36 if len(num) <= 4 else 26
    tb(s10, num,    vx+Inches(0.2), VY+Inches(0.12), VW-Inches(0.3), Inches(0.65),
       size=nsize, bold=True, color=ORANGE)
    tb(s10, label,  vx+Inches(0.2), VY+Inches(0.8),  VW-Inches(0.3), Inches(0.65),
       size=10, color=MUTED)
    tb(s10, source, vx+Inches(0.2), VY+VH-Inches(0.32), VW-Inches(0.3), Inches(0.3),
       size=9, color=FOOTNOTE)

rect(s10, Inches(0.6), Inches(3.95), Inches(12.1), Inches(2.72), CARD_BG)
rect(s10, Inches(0.6), Inches(3.95), Inches(12.1), Inches(0.05), ORANGE)
tb(s10, 'USER RESEARCH  —  12 INTERVIEWS  ·  LAGOS & ABUJA',
   Inches(0.85), Inches(4.07), Inches(11.0), Inches(0.38), size=11, bold=True, color=ORANGE)

for i, finding in enumerate([
    '12 of 12  —  had a smartphone and used mobile money (OPay, PalmPay) daily. Digital exclusion is not the problem.',
    '9 of 12   —  had been rejected for a bank loan or could not access credit in the past 12 months.',
    '11 of 12  —  said they would share financial data in exchange for better access to loans and financial products.',
    '8 of 12   —  were earning across multiple income streams with no institution able to see the full picture.',
    '"Every person we spoke to knew exactly what they earned and what they were worth. The system just couldn\'t see it."',
]):
    is_quote = i == 4
    tb(s10, finding, Inches(0.85), Inches(4.53 + i * 0.42), Inches(11.5), Inches(0.4),
       size=12, color=ORANGE if is_quote else MUTED, italic=is_quote)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — THE TEAM
# ══════════════════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(blank)
base(s11, '10  —  THE TEAM')
tb(s11, 'The Right People for This Moment.',
   Inches(0.6), Inches(0.85), Inches(10), Inches(0.75), size=38, bold=True)
tb(s11,
   'We are Nigerian engineers who built Trace because this problem is personal. '
   'We know these users — because we grew up around them.',
   Inches(0.6), Inches(1.6), Inches(11.5), Inches(0.5), size=13, color=MUTED)

TW = Inches(3.75)
TH = Inches(3.2)
TY = Inches(2.25)

for (tx, name, role, skills, tagline) in [
    (Inches(0.6),
     'Praise Olaiwon',
     'Product Lead & Full-Stack Engineer',
     'System architecture  ·  NestJS backend  ·  Next.js frontend\n'
     'Squad API integration  ·  PostgreSQL  ·  Docker  ·  CI/CD pipeline',
     '"Built the full platform end-to-end"'),
    (Inches(4.75),
     '[Team Member 2]',
     'AI / ML Engineer',
     'Python FastAPI  ·  gRPC scoring service  ·  scikit-learn\n'
     'Economic identity model  ·  job matching algorithm  ·  BullMQ',
     '"Built the intelligence layer that makes Trace smarter"'),
    (Inches(8.9),
     '[Team Member 3]',
     'Frontend & Product Designer',
     'React / Next.js  ·  Mobile-first PWA  ·  Onboarding UX\n'
     'Dashboard design  ·  User research interviews  ·  Figma',
     '"Designed the experience from the user\'s perspective"'),
]:
    card_base(s11, tx, TY, TW, TH)
    tb(s11, name,    tx+Inches(0.25), TY+Inches(0.18), TW-Inches(0.4), Inches(0.5),
       size=15, bold=True)
    tb(s11, role,    tx+Inches(0.25), TY+Inches(0.68), TW-Inches(0.4), Inches(0.38),
       size=11, color=ORANGE)
    rect(s11, tx+Inches(0.25), TY+Inches(1.1), TW-Inches(0.4), Inches(0.015), DIM)
    tb(s11, skills,  tx+Inches(0.25), TY+Inches(1.18), TW-Inches(0.4), Inches(1.45),
       size=11, color=MUTED)
    tb(s11, tagline, tx+Inches(0.25), TY+TH-Inches(0.62), TW-Inches(0.4), Inches(0.55),
       size=11, italic=True, color=ORANGE)

rect(s11, Inches(0.6), Inches(5.65), Inches(12.1), Inches(1.2), RGBColor(0x06, 0x0e, 0x1f))
rect(s11, Inches(0.6), Inches(5.65), Inches(0.05),  Inches(1.2), ORANGE)
tb(s11,
   '"We didn\'t build Trace to win a hackathon. We built it because 39 million Nigerians do real '
   'work every day and the financial system doesn\'t know they exist. '
   'Trace is how we fix that — one verified identity at a time."',
   Inches(0.85), Inches(5.8), Inches(11.6), Inches(0.95),
   size=14, italic=True)


# ── Save ──────────────────────────────────────────────────────────────────────
out = r'c:\Users\USER\Desktop\squad-hackathon\docs\Trace_Pitch_Deck.pptx'
prs.save(out)
print(f'Saved: {out}')
